#!/usr/bin/env python3
"""Extract documents to Markdown + JSON for Loci."""

from __future__ import annotations

import argparse
import json
import os
import platform
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any

MIN_OK_WORDS = 20

LIBREOFFICE_CANDIDATES = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/Applications/LibreOffice.app/Contents/MacOS/libreoffice",
    "/usr/local/bin/soffice",
    "/opt/homebrew/bin/soffice",
]


def configure_environment() -> None:
    if os.environ.get("DOCLING_LIBREOFFICE_CMD"):
        return
    for candidate in LIBREOFFICE_CANDIDATES:
        if Path(candidate).exists():
            os.environ["DOCLING_LIBREOFFICE_CMD"] = candidate
            return
    discovered = shutil.which("soffice")
    if discovered:
        os.environ["DOCLING_LIBREOFFICE_CMD"] = discovered


def word_count(text: str) -> int:
    return len(re.findall(r"\S+", text))


def write_outputs(output_dir: Path, markdown: str, document: Any | None, meta: dict[str, Any]) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    markdown = markdown.strip()
    (output_dir / "extracted.md").write_text(markdown + "\n", encoding="utf-8")
    (output_dir / "extracted.txt").write_text(markdown + "\n", encoding="utf-8")
    if document is not None:
        (output_dir / "document.json").write_text(
            json.dumps(document, indent=2, ensure_ascii=False) + "\n",
            encoding="utf-8",
        )
    meta["word_count"] = word_count(markdown)
    meta["status"] = "ok" if meta["word_count"] >= MIN_OK_WORDS else "thin"
    (output_dir / "extract-meta.json").write_text(
        json.dumps(meta, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )


def try_docling(input_path: Path) -> tuple[str | None, Any | None, dict[str, Any]]:
    configure_environment()
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        return None, None, {}

    meta: dict[str, Any] = {
        "extractor": "docling",
        "ocr_used": False,
        "warnings": [],
        "libreoffice": os.environ.get("DOCLING_LIBREOFFICE_CMD"),
    }
    try:
        converter = DocumentConverter()
        result = converter.convert(str(input_path))
        document = result.document
        markdown = document.export_to_markdown()
        document_json: Any | None = None
        if hasattr(document, "export_to_dict"):
            document_json = document.export_to_dict()
        elif hasattr(document, "export_to_document_dict"):
            document_json = document.export_to_document_dict()
        return markdown, document_json, meta
    except Exception as exc:  # noqa: BLE001
        meta["warnings"].append(f"docling failed: {exc}")
        return None, None, meta


def textutil_to_text(input_path: Path) -> str | None:
    if platform.system() != "Darwin" or not shutil.which("textutil"):
        return None
    try:
        completed = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(input_path)],
            check=True,
            capture_output=True,
            text=True,
        )
        text = completed.stdout.strip()
        return text or None
    except subprocess.CalledProcessError:
        return None


def pypdf_text(input_path: Path) -> str | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        return None
    try:
        reader = PdfReader(str(input_path))
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(pages).strip()
        return text or None
    except Exception:
        return None


def docx_text(input_path: Path) -> str | None:
    try:
        from docx import Document
    except ImportError:
        return None
    try:
        document = Document(str(input_path))
        paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs).strip() or None
    except Exception:
        return None


def pptx_text(input_path: Path) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        presentation = Presentation(str(input_path))
        chunks: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text.strip())
            if lines:
                chunks.append(f"## Slide {index}\n\n" + "\n".join(lines))
        return "\n\n".join(chunks).strip() or None
    except Exception:
        return None


def xlsx_text(input_path: Path) -> str | None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    try:
        workbook = load_workbook(str(input_path), read_only=True, data_only=True)
        chunks: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                chunks.append(f"## {sheet.title}\n\n" + "\n".join(rows))
        return "\n\n".join(chunks).strip() or None
    except Exception:
        return None


def image_text(input_path: Path) -> tuple[str | None, bool]:
    try:
        from PIL import Image
    except ImportError:
        return None, False
    try:
        import pytesseract
    except ImportError:
        return None, False
    try:
        image = Image.open(input_path)
        text = pytesseract.image_to_string(image).strip()
        return (text or None), True
    except Exception:
        return None, True


def plain_text(input_path: Path) -> str | None:
    try:
        return input_path.read_text(encoding="utf-8").strip() or None
    except UnicodeDecodeError:
        return None


def fallback_extract(input_path: Path) -> tuple[str | None, Any | None, dict[str, Any]]:
    ext = input_path.suffix.lower()
    meta: dict[str, Any] = {
        "extractor": "fallback",
        "ocr_used": False,
        "warnings": [],
    }
    text: str | None = None

    if ext in {".txt", ".md", ".markdown", ".html", ".htm", ".json", ".csv", ".xml"}:
        text = plain_text(input_path)
        meta["extractor"] = "plain-text"
    elif ext == ".pdf":
        text = pypdf_text(input_path)
        meta["extractor"] = "pypdf"
    elif ext in {".doc", ".docx", ".rtf", ".html", ".htm"}:
        text = docx_text(input_path) or textutil_to_text(input_path)
        meta["extractor"] = "python-docx" if ext == ".docx" and text else "textutil"
    elif ext in {".ppt", ".pptx"}:
        text = pptx_text(input_path) or textutil_to_text(input_path)
        meta["extractor"] = "python-pptx" if ext == ".pptx" and text else "textutil"
    elif ext in {".xls", ".xlsx"}:
        text = xlsx_text(input_path)
        meta["extractor"] = "openpyxl"
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".tif", ".tiff", ".bmp", ".heic"}:
        text, ocr_used = image_text(input_path)
        meta["extractor"] = "tesseract"
        meta["ocr_used"] = ocr_used
        if not text and ocr_used:
            meta["warnings"].append("OCR produced no text")
        if not ocr_used:
            meta["warnings"].append("image OCR unavailable (install pillow + pytesseract or docling)")
    else:
        text = textutil_to_text(input_path)
        meta["extractor"] = "textutil"

    document_json = {
        "source": str(input_path),
        "extractor": meta["extractor"],
        "format": ext.lstrip("."),
        "blocks": [{"type": "text", "text": text}] if text else [],
    }
    return text, document_json, meta


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract a document for Loci")
    parser.add_argument("--input", required=True, help="Path to source file")
    parser.add_argument("--output-dir", required=True, help="Directory for extracted artifacts")
    args = parser.parse_args()

    input_path = Path(args.input).expanduser().resolve()
    output_dir = Path(args.output_dir).expanduser().resolve()

    if not input_path.exists():
        print(f"input not found: {input_path}", file=sys.stderr)
        return 2

    meta: dict[str, Any] = {
        "input": str(input_path),
        "extractor": None,
        "ocr_used": False,
        "word_count": 0,
        "status": "failed",
        "warnings": [],
        "libreoffice": os.environ.get("DOCLING_LIBREOFFICE_CMD"),
    }

    markdown, document_json, extractor_meta = try_docling(input_path)
    if markdown:
        meta.update(extractor_meta)
    else:
        if extractor_meta.get("warnings"):
            meta["warnings"].extend(extractor_meta["warnings"])
        markdown, document_json, extractor_meta = fallback_extract(input_path)
        meta.update({k: v for k, v in extractor_meta.items() if k != "warnings"})
        if extractor_meta.get("warnings"):
            meta["warnings"].extend(extractor_meta["warnings"])

    if markdown and markdown.strip():
        write_outputs(output_dir, markdown, document_json, meta)
        return 0

    meta["status"] = "failed"
    meta["warnings"].append("no extractor produced usable text")
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "extract-meta.json").write_text(
        json.dumps(meta, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    return 1


if __name__ == "__main__":
    raise SystemExit(main())
